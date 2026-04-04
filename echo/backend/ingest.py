"""
ingest.py — Lit les fichiers d'un dossier (PDF, DOCX, PPTX, XLSX, Markdown, TXT),
les découpe en chunks, les vectorise via Ollama et les stocke dans ChromaDB.

Formats supportés :
  - .pdf         → pdfplumber (layout avancé) + fallback pypdf
  - .docx        → python-docx (paragraphes + tableaux)
  - .pptx        → python-pptx (texte des slides + notes)
  - .xlsx / .xls → openpyxl (cellules de toutes les feuilles)
  - .md / .txt   → lecture UTF-8 directe
"""

import os
import hashlib
import logging
from pathlib import Path
from typing import Generator

import chromadb
import requests

import config

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# ChromaDB client
# ---------------------------------------------------------------------------

def get_collection() -> chromadb.Collection:
    client = chromadb.PersistentClient(path=config.CHROMA_DIR)
    return client.get_or_create_collection(
        name=config.COLLECTION_NAME,
        metadata={"hnsw:space": "cosine"},
    )


# ---------------------------------------------------------------------------
# Text extraction — PDF
# ---------------------------------------------------------------------------

def extract_pdf(path: Path) -> str:
    """
    Extraction PDF en deux passes :
    1. pdfplumber — meilleure gestion des colonnes, tableaux et espaces
    2. pypdf      — fallback léger si pdfplumber échoue ou produit trop peu de texte

    Les PDFs purement scannés (image sans couche texte) produiront un résultat
    vide : c'est une limitation connue, affichée en WARNING.
    """
    text_parts: list[str] = []

    # — Passe 1 : pdfplumber —
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                try:
                    t = page.extract_text(x_tolerance=2, y_tolerance=2)
                    if t:
                        text_parts.append(t)
                except Exception as page_err:
                    log.warning("pdfplumber — page ignorée dans %s : %s", path.name, page_err)
        if text_parts:
            return "\n\n".join(text_parts)
    except ImportError:
        log.debug("pdfplumber non installé — passage au fallback pypdf")
    except Exception as err:
        log.warning("pdfplumber échoué sur %s (%s) — tentative avec pypdf", path.name, err)

    # — Passe 2 : pypdf (fallback) —
    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        for page in reader.pages:
            try:
                t = page.extract_text() or ""
                if t.strip():
                    text_parts.append(t)
            except Exception as page_err:
                log.warning("pypdf — page ignorée dans %s : %s", path.name, page_err)
        if text_parts:
            return "\n\n".join(text_parts)
    except ImportError:
        log.warning("pypdf non installé — impossible de lire %s", path.name)
    except Exception as err:
        log.warning("pypdf échoué sur %s : %s", path.name, err)

    log.warning(
        "Aucun texte extrait de %s — le fichier est peut-être un PDF scanné (image uniquement).",
        path.name,
    )
    return ""


# ---------------------------------------------------------------------------
# Text extraction — DOCX (Word)
# ---------------------------------------------------------------------------

def extract_docx(path: Path) -> str:
    """
    Extrait le texte d'un fichier Word (.docx) :
    - Tous les paragraphes dans l'ordre du document
    - Le contenu des tableaux (cellule par cellule)
    """
    try:
        from docx import Document
    except ImportError:
        log.warning("python-docx non installé — skipping %s", path.name)
        return ""

    try:
        doc = Document(str(path))
        parts: list[str] = []

        for element in doc.element.body:
            tag = element.tag.split("}")[-1]  # retire le namespace XML

            if tag == "p":
                # Paragraphe
                from docx.oxml.ns import qn
                text = "".join(
                    node.text or ""
                    for node in element.iter(qn("w:t"))
                )
                if text.strip():
                    parts.append(text)

            elif tag == "tbl":
                # Tableau — on lit chaque cellule
                from docx.oxml.ns import qn
                rows: list[str] = []
                for row in element.iter(qn("w:tr")):
                    cells = [
                        "".join(node.text or "" for node in cell.iter(qn("w:t")))
                        for cell in row.iter(qn("w:tc"))
                    ]
                    row_text = " | ".join(c.strip() for c in cells if c.strip())
                    if row_text:
                        rows.append(row_text)
                if rows:
                    parts.append("\n".join(rows))

        return "\n\n".join(parts)

    except Exception as err:
        log.warning("Erreur lecture DOCX %s : %s", path.name, err)
        return ""


# ---------------------------------------------------------------------------
# Text extraction — PPTX (PowerPoint)
# ---------------------------------------------------------------------------

def extract_pptx(path: Path) -> str:
    """
    Extrait le texte d'un fichier PowerPoint (.pptx) slide par slide :
    - Titre + corps de chaque diapositive
    - Notes du présentateur (si présentes)
    """
    try:
        from pptx import Presentation
    except ImportError:
        log.warning("python-pptx non installé — skipping %s", path.name)
        return ""

    try:
        prs = Presentation(str(path))
        slides_text: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            parts: list[str] = []

            # Texte des shapes (title, content, text boxes…)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)

            # Notes du présentateur
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes = notes_tf.text.strip()
                if notes:
                    parts.append(f"[Notes slide {i}] {notes}")

            if parts:
                slides_text.append(f"— Slide {i} —\n" + "\n".join(parts))

        return "\n\n".join(slides_text)

    except Exception as err:
        log.warning("Erreur lecture PPTX %s : %s", path.name, err)
        return ""


# ---------------------------------------------------------------------------
# Text extraction — XLSX (Excel)
# ---------------------------------------------------------------------------

def extract_xlsx(path: Path) -> str:
    """
    Extrait le texte d'un fichier Excel (.xlsx / .xls) :
    - Toutes les feuilles, toutes les lignes non vides
    - Format : « Feuille | col1 | col2 | … »
    """
    try:
        import openpyxl
    except ImportError:
        log.warning("openpyxl non installé — skipping %s", path.name)
        return ""

    try:
        # read_only=True pour les grands fichiers, data_only=True pour les valeurs calculées
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"=== Feuille : {sheet_name} ===\n" + "\n".join(rows))

        wb.close()
        return "\n\n".join(parts)

    except Exception as err:
        log.warning("Erreur lecture XLSX %s : %s", path.name, err)
        return ""


# ---------------------------------------------------------------------------
# Text extraction — Markdown / TXT
# ---------------------------------------------------------------------------

def extract_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------

EXTRACTORS = {
    # Texte brut
    ".md":       extract_text_file,
    ".markdown": extract_text_file,
    ".txt":      extract_text_file,
    # PDF
    ".pdf":      extract_pdf,
    # Office
    ".docx":     extract_docx,
    ".pptx":     extract_pptx,
    ".xlsx":     extract_xlsx,
    ".xls":      extract_xlsx,
}


def extract_text(path: Path) -> str:
    extractor = EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        return ""
    return extractor(path)


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_text(
    text: str,
    size: int = config.CHUNK_SIZE,
    overlap: int = config.CHUNK_OVERLAP,
) -> Generator[str, None, None]:
    """Sliding-window character chunker avec respect des fins de phrase."""
    text = text.strip()
    if not text:
        return
    start = 0
    while start < len(text):
        end = start + size
        # Tente de couper sur une fin de phrase plutôt qu'en plein milieu d'un mot
        if end < len(text):
            for sep in ("\n\n", "\n", ". ", " "):
                pos = text.rfind(sep, start + size // 2, end)
                if pos != -1:
                    end = pos + len(sep)
                    break
        yield text[start:end]
        start += size - overlap


# ---------------------------------------------------------------------------
# Embedding via Ollama
# ---------------------------------------------------------------------------

def embed(texts: list[str]) -> list[list[float]]:
    """Call Ollama /api/embed for a batch of texts."""
    url = f"{config.OLLAMA_BASE_URL}/api/embed"
    response = requests.post(
        url,
        json={"model": config.OLLAMA_EMBED_MODEL, "input": texts},
        timeout=120,
    )
    if response.status_code == 404:
        raise RuntimeError(
            f"Modèle d'embedding '{config.OLLAMA_EMBED_MODEL}' introuvable dans Ollama. "
            f"Installez-le avec : ollama pull {config.OLLAMA_EMBED_MODEL}"
        )
    response.raise_for_status()
    return response.json()["embeddings"]


# ---------------------------------------------------------------------------
# Ingestion pipeline
# ---------------------------------------------------------------------------

def doc_id(file_path: Path, chunk_index: int) -> str:
    base = hashlib.md5(str(file_path.resolve()).encode()).hexdigest()[:8]
    return f"{base}_{chunk_index}"


def ingest_file(path: Path, collection: chromadb.Collection) -> int:
    text = extract_text(path)
    if not text.strip():
        log.warning(
            "Aucun texte extrait de %s — fichier vide, scanné (image PDF) ou format non supporté.",
            path.name,
        )
        return 0

    chunks = list(chunk_text(text))
    if not chunks:
        return 0

    embeddings = embed(chunks)

    ids = [doc_id(path, i) for i in range(len(chunks))]
    mtime = path.stat().st_mtime
    metadatas = [
        {"source": str(path), "filename": path.name, "chunk": i, "mtime": mtime}
        for i in range(len(chunks))
    ]

    collection.upsert(
        ids=ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metadatas,
    )
    log.info("  ✓ %s — %d chunks ingérés", path.name, len(chunks))
    return len(chunks)


def ingest_directory(docs_dir: str = config.DOCS_DIR) -> None:
    docs_path = Path(docs_dir)
    if not docs_path.exists():
        log.error("Dossier introuvable : %s", docs_dir)
        return

    collection = get_collection()
    total = 0

    files = [p for p in docs_path.rglob("*") if p.suffix.lower() in EXTRACTORS]
    if not files:
        log.warning(
            "Aucun fichier supporté trouvé dans %s (PDF, DOCX, PPTX, XLSX, MD, TXT)",
            docs_dir,
        )
        return

    log.info("Ingestion de %d fichier(s) depuis %s", len(files), docs_dir)
    for file in files:
        total += ingest_file(file, collection)

    log.info("Ingestion terminée — %d chunks au total dans ChromaDB", total)


if __name__ == "__main__":
    ingest_directory()
